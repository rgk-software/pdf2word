from pptx import Presentation
from docx import Document

def create_pptx(filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"
    prs.save(filename)
    print(f"Created {filename}")

def create_docx(filename):
    doc = Document()
    doc.add_heading('Document Title', 0)
    doc.add_paragraph('A plain paragraph having some bold and some italic.')
    doc.save(filename)
    print(f"Created {filename}")

if __name__ == "__main__":
    create_pptx("test.pptx")
    create_docx("test.docx")
