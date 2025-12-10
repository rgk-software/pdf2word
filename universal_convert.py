import sys
import os
import zipfile
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_docx(input_file, output_file):
    cv = Converter(input_file)
    cv.convert(output_file)
    cv.close()

def pdf_to_pptx(input_file, output_file):
    prs = Presentation()
    # Use blank layout (usually index 6)
    blank_slide_layout = prs.slide_layouts[6]
    
    images = convert_from_path(input_file)
    
    # Calculate aspect ratio of the first image to set slide size?
    # For now, default 4:3 or 16:9 
    # Let's resize slides to match image? 
    if images:
        width, height = images[0].size
        # Set presentation slide size (emu = pixels * 9525 ?)
        # PPTX measures in EMUs. 1 inch = 914400 EMUs.
        # approx 72 dpi or 96 dpi for pixels? 
        # Easier: just add image to standard slide and center/fit.
        
        # Actually, best experience: resize slide matching PDF page.
        prs.slide_width = int(width * 9525) 
        prs.slide_height = int(height * 9525)

    for i, image in enumerate(images):
        slide = prs.slides.add_slide(blank_slide_layout)
        image_path = f"temp_slide_{i}.png"
        image.save(image_path)
        
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        os.remove(image_path)
        
    prs.save(output_file)

def pdf_to_image(input_file, output_file):
    images = convert_from_path(input_file)
    if len(images) == 1:
        images[0].save(output_file, 'PNG')
    else:
        # Zip multiple images
        base_name = os.path.splitext(output_file)[0]
        zip_filename = base_name + ".zip"
        
        with zipfile.ZipFile(zip_filename, 'w') as zipf:
            for i, image in enumerate(images):
                img_name = f"page_{i+1}.png"
                image.save(img_name, 'PNG')
                zipf.write(img_name)
                os.remove(img_name)
        
        # Rename zip to the expected output filename if the caller expects it?
        # Caller expects 'output_file'. If we made a zip, we should probably output that.
        # But for 'if single page', the requirement is image.
        # We will handle the zip vs png logic here.
        # If we created a zip, we replace output_file with it.
        if output_file.endswith('.png') or output_file.endswith('.jpg'):
             # We can't save a zip as .png. 
             # We rely on the backend to name it correctly?
             # Or we rename the zip to the output path.
             os.rename(zip_filename, output_file)

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python universal_convert.py <mode> <input> <output>")
        sys.exit(1)
    
    mode = sys.argv[1]
    input_file = sys.argv[2]
    output_file = sys.argv[3]
    
    try:
        if mode == 'pdf-to-word':
            pdf_to_docx(input_file, output_file)
        elif mode == 'pdf-to-ppt':
            pdf_to_pptx(input_file, output_file)
        elif mode == 'pdf-to-image':
            pdf_to_image(input_file, output_file)
        else:
            print(f"Unknown mode: {mode}")
            sys.exit(1)
        print("Conversion Successful")
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
